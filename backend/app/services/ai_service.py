import json
from pathlib import Path
from typing import Any

from sqlalchemy.orm import Session

from app.core.security import get_password_hash
from app.models.entities import (
    AIReportHistory,
    Conversation,
    Document,
    FileResource,
    Message,
    Milestone,
    Notification,
    Project,
    ProjectMember,
    RiskAlert,
    Task,
    TaskActivity,
    User,
)
from app.services.llm_service import LLMServiceError, complete_json, complete_text
from app.services.member_insights import build_contribution_payload
from app.services.presenters import (
    build_dashboard_payload,
    build_git_payload,
    build_graph_payload,
    build_reminder_payload,
    build_user_map,
    serialize_conversation,
    serialize_message,
    serialize_task,
)
from app.utils.serializers import dump_json, parse_json


AI_ASSISTANT_USERNAME = "teamflow_ai"
AI_ASSISTANT_NAME = "TeamFlow AI"

ALLOWED_ACTIONS = {
    "navigate_page",
    "open_task",
    "open_document",
    "open_file",
    "create_task",
    "update_task",
    "accept_task",
    "block_task",
    "complete_task",
    "create_milestone",
    "generate_weekly_report",
    "scan_risks",
    "recalculate_contribution",
}


def get_or_create_ai_user(db: Session) -> User:
    row = db.query(User).filter(User.username == AI_ASSISTANT_USERNAME).first()
    if row:
        return row
    row = User(
        username=AI_ASSISTANT_USERNAME,
        password_hash=get_password_hash("teamflow-ai-system"),
        display_name=AI_ASSISTANT_NAME,
        email="teamflow-ai@local",
        avatar="https://api.dicebear.com/7.x/bottts/svg?seed=teamflow-ai&backgroundColor=0f172a,1d4ed8,0f766e",
        system_role="system",
        title="项目智能助手",
        bio="负责项目上下文分析、实时建议与安全范围内的自动化执行。",
    )
    db.add(row)
    db.commit()
    db.refresh(row)
    return row


def get_or_create_ai_conversation(db: Session, project_id: int, created_by: int) -> Conversation:
    row = (
        db.query(Conversation)
        .filter(Conversation.project_id == project_id, Conversation.conversation_type == "ai_assistant")
        .first()
    )
    if row:
        return row
    row = Conversation(
        project_id=project_id,
        conversation_type="ai_assistant",
        name="TeamFlow AI 助手",
        created_by=created_by,
        is_pinned=True,
    )
    db.add(row)
    db.commit()
    db.refresh(row)
    return row


def persist_message(
    db: Session,
    *,
    project_id: int,
    conversation_id: int,
    sender_id: int,
    content: str,
    message_type: str,
    metadata: dict[str, Any] | None = None,
    code_language: str = "",
) -> Message:
    row = Message(
        project_id=project_id,
        conversation_id=conversation_id,
        sender_id=sender_id,
        message_type=message_type,
        content=content,
        code_language=code_language,
        metadata_json=dump_json(metadata or {}),
    )
    db.add(row)
    db.commit()
    db.refresh(row)
    return row


def list_ai_conversation_payload(db: Session, project_id: int, current_user_id: int) -> dict[str, Any]:
    conversation = get_or_create_ai_conversation(db, project_id, current_user_id)
    users = build_user_map(db)
    messages = (
        db.query(Message)
        .filter(Message.conversation_id == conversation.id)
        .order_by(Message.created_at.asc())
        .all()
    )
    return {
        "conversation": serialize_conversation(conversation, db, current_user_id),
        "messages": [serialize_message(item, users) for item in messages],
    }


def build_project_snapshot(
    db: Session,
    project_id: int,
    *,
    route_name: str = "",
    route_path: str = "",
    route_params: dict[str, Any] | None = None,
    query: dict[str, Any] | None = None,
    page_context: dict[str, Any] | None = None,
) -> dict[str, Any]:
    project = db.query(Project).filter(Project.id == project_id).first()
    members = db.query(ProjectMember).filter(ProjectMember.project_id == project_id).all()
    users = build_user_map(db)
    tasks = db.query(Task).filter(Task.project_id == project_id, Task.archived == False).order_by(Task.updated_at.desc()).all()  # noqa: E712
    documents = db.query(Document).filter(Document.project_id == project_id, Document.is_deleted == False).order_by(Document.updated_at.desc()).all()  # noqa: E712
    files = db.query(FileResource).filter(FileResource.project_id == project_id).order_by(FileResource.updated_at.desc()).all()
    route_params = route_params or {}
    query = query or {}
    page_context = page_context or {}

    selected_task_id = _coerce_int(page_context.get("selectedTaskId") or query.get("taskId") or route_params.get("taskId"))
    selected_document_id = _coerce_int(page_context.get("selectedDocumentId") or query.get("docId") or route_params.get("docId"))
    selected_file_id = _coerce_int(page_context.get("selectedFileId") or query.get("fileId"))
    selected_conversation_id = _coerce_int(page_context.get("selectedConversationId") or query.get("conversationId"))
    selected_user_id = _coerce_int(page_context.get("selectedUserId") or query.get("userId"))

    page_data: dict[str, Any]
    if route_name == "dashboard":
        page_data = build_dashboard_payload(project_id, db)
    elif route_name == "tasks":
        page_data = {
            "tasks": [_task_brief(item, users) for item in tasks[:12]],
            "selected_task": _task_brief(next((item for item in tasks if item.id == selected_task_id), None), users),
            "filters": page_context.get("filters", {}),
        }
    elif route_name == "messages":
        page_data = _messages_page_payload(db, project_id, selected_conversation_id)
    elif route_name == "documents":
        page_data = {
            "documents": [_document_brief(item) for item in documents[:10]],
            "selected_document": _document_brief(next((item for item in documents if item.id == selected_document_id), None), include_content=True),
        }
    elif route_name == "document-editor":
        page_data = {
            "selected_document": _document_brief(next((item for item in documents if item.id == selected_document_id), None), include_content=True),
        }
    elif route_name == "files":
        page_data = {
            "files": [_file_brief(item) for item in files[:10]],
            "selected_file": _file_brief(next((item for item in files if item.id == selected_file_id), None), include_description=True),
            "filters": page_context.get("filters", {}),
        }
    elif route_name == "graph":
        page_data = build_graph_payload(project_id, db)
    elif route_name == "contribution":
        page_data = build_contribution_payload(db, project_id, prefer_llm=False, persist=False)
    elif route_name == "career":
        contribution = build_contribution_payload(db, project_id, prefer_llm=False, persist=False)
        profiles = contribution.get("member_profiles", [])
        page_data = {
            "summary": contribution.get("summary", {}),
            "selected_profile": next((item for item in profiles if item["user"]["id"] == selected_user_id), profiles[0] if profiles else None),
        }
    elif route_name == "risk":
        page_data = {"risks": [parse_json(dump_json(item), {}) for item in build_dashboard_payload(project_id, db).get("risk_alerts", [])]}
    elif route_name == "git":
        page_data = build_git_payload(project_id, db)
    elif route_name == "settings":
        page_data = {
            "project": {
                "name": project.name if project else "",
                "course_name": project.course_name if project else "",
                "description": project.description if project else "",
                "due_date": project.due_date if project else "",
            },
            "members": [
                {"id": row.user_id, "name": users[row.user_id].display_name, "role": row.role}
                for row in members
                if row.user_id in users
            ],
        }
    elif route_name == "reminders":
        page_data = build_reminder_payload(project_id, _coerce_int(page_context.get("currentUserId")) or 0, db)
    elif route_name == "milestones":
        milestones = db.query(Milestone).filter(Milestone.project_id == project_id).order_by(Milestone.order_index.asc()).all()
        page_data = {
            "milestones": [
                {"id": item.id, "name": item.name, "due_date": item.due_date, "status": item.status}
                for item in milestones
            ]
        }
    elif route_name == "ai":
        page_data = {
            "recent_reports": [
                {
                    "id": item.id,
                    "report_type": item.report_type,
                    "created_at": str(item.created_at),
                    "content": item.content[:500],
                }
                for item in db.query(AIReportHistory)
                .filter(AIReportHistory.project_id == project_id)
                .order_by(AIReportHistory.created_at.desc())
                .limit(5)
                .all()
            ],
            "task_count": len(tasks),
            "document_count": len(documents),
        }
    else:
        page_data = {
            "task_count": len(tasks),
            "document_count": len(documents),
            "file_count": len(files),
        }

    return {
        "project": {
            "id": project.id if project else project_id,
            "name": project.name if project else "",
            "course_name": project.course_name if project else "",
            "description": project.description if project else "",
            "due_date": project.due_date if project else "",
            "repo_url": project.repo_url if project else "",
            "member_count": len(members),
            "task_count": len(tasks),
            "document_count": len(documents),
            "file_count": len(files),
        },
        "members": [
            {"id": row.user_id, "name": users[row.user_id].display_name, "role": row.role}
            for row in members
            if row.user_id in users
        ],
        "current_page": {
            "route_name": route_name,
            "route_path": route_path,
            "route_params": route_params,
            "query": query,
            "page_context": page_context,
        },
        "highlights": {
            "open_tasks": len([item for item in tasks if item.status != "DONE"]),
            "blocked_tasks": len([item for item in tasks if item.status == "BLOCKED"]),
            "done_tasks": len([item for item in tasks if item.status == "DONE"]),
        },
        "page_data": page_data,
    }


def generate_project_plan(db: Session, project_id: int, user_prompt: str) -> dict[str, Any]:
    snapshot = build_project_snapshot(db, project_id, route_name="ai")
    messages = [
        {
            "role": "system",
            "content": (
                "你是 TeamFlow 项目智能助手。请基于真实项目上下文输出严格 JSON，"
                "字段必须包含 summary, phases, tasks, risks, suggestions。"
                "phases 为数组，每项含 step/title/date；"
                "tasks 为数组，每项含 name/owner/priority/hours/deadline/status；"
                "risks 为数组，每项含 name/level/suggestion；"
                "suggestions 为字符串数组。不要输出 Markdown。"
            ),
        },
        {
            "role": "user",
            "content": json.dumps({"user_prompt": user_prompt, "project_snapshot": snapshot}, ensure_ascii=False),
        },
    ]
    parsed = complete_json(messages, temperature=0.25, max_tokens=1800, retries=1)
    parsed.setdefault("summary", f"围绕 {snapshot['project']['name']} 生成的实时任务规划。")
    parsed.setdefault("phases", [])
    parsed.setdefault("tasks", [])
    parsed.setdefault("risks", [])
    parsed.setdefault("suggestions", [])
    return parsed


def generate_weekly_report(db: Session, project_id: int) -> str:
    snapshot = build_project_snapshot(db, project_id, route_name="dashboard")
    messages = [
        {
            "role": "system",
            "content": (
                "你是 TeamFlow 周报助手。请基于真实项目数据生成中文周报，"
                "输出纯文本，包含本周进展、风险、下周建议三部分。"
            ),
        },
        {
            "role": "user",
            "content": json.dumps(snapshot, ensure_ascii=False),
        },
    ]
    return complete_text(messages, temperature=0.35, max_tokens=1400)


def analyze_document(db: Session, *, project_id: int, document_id: int, prompt: str = "") -> dict[str, Any]:
    document = (
        db.query(Document)
        .filter(Document.project_id == project_id, Document.id == document_id, Document.is_deleted == False)  # noqa: E712
        .first()
    )
    if not document:
        raise ValueError("文档不存在")
    return _analyze_content(
        source_name=document.title,
        source_type="document",
        content=document.content,
        prompt=prompt,
    )


def analyze_file(db: Session, *, project_id: int, file_id: int, prompt: str = "") -> dict[str, Any]:
    file_row = db.query(FileResource).filter(FileResource.project_id == project_id, FileResource.id == file_id).first()
    if not file_row:
        raise ValueError("文件不存在")
    content = extract_file_text(file_row.storage_path, file_row.name)
    return _analyze_content(
        source_name=file_row.name,
        source_type="file",
        content=content,
        prompt=prompt,
    )


def _analyze_content(*, source_name: str, source_type: str, content: str, prompt: str) -> dict[str, Any]:
    excerpt = (content or "").strip()
    if not excerpt:
        raise ValueError("文件内容为空，无法分析")
    excerpt = excerpt[:12000]
    messages = [
        {
            "role": "system",
            "content": (
                "你是 TeamFlow 文档分析助手。请输出严格 JSON，字段包含 "
                "summary, scale, workload, grade, risks, suggestions, extracted_excerpt。"
                "risks 和 suggestions 必须是字符串数组。不要输出 Markdown。"
            ),
        },
        {
            "role": "user",
            "content": json.dumps(
                {
                    "source_name": source_name,
                    "source_type": source_type,
                    "analysis_goal": prompt or "请分析文档质量、工作量、风险点与改进建议。",
                    "content_excerpt": excerpt,
                },
                ensure_ascii=False,
            ),
        },
    ]
    parsed = complete_json(messages, temperature=0.2, max_tokens=1400, retries=1)
    parsed.setdefault("summary", "")
    parsed.setdefault("scale", f"{len(excerpt)} 字")
    parsed.setdefault("workload", "中")
    parsed.setdefault("grade", "A")
    parsed.setdefault("risks", [])
    parsed.setdefault("suggestions", [])
    parsed.setdefault("extracted_excerpt", excerpt[:300])
    parsed["source_name"] = source_name
    parsed["source_type"] = source_type
    return parsed


def build_chat_completion(
    db: Session,
    *,
    project_id: int,
    prompt: str,
    route_name: str,
    route_path: str,
    route_params: dict[str, Any] | None,
    query: dict[str, Any] | None,
    page_context: dict[str, Any] | None,
) -> dict[str, Any]:
    snapshot = build_project_snapshot(
        db,
        project_id,
        route_name=route_name,
        route_path=route_path,
        route_params=route_params,
        query=query,
        page_context=page_context,
    )
    messages = [
        {
            "role": "system",
            "content": (
                "你是 TeamFlow 的项目 AI 助手。你可以结合当前页面上下文回答问题，并在用户意图明确时自动安排非破坏性动作。"
                "请严格输出 JSON，包含 reply 和 actions 两个字段。"
                "reply 为给用户看的中文回答；actions 为数组，每项包含 type, params, reason, ui_target。"
                f"只允许使用这些动作: {sorted(ALLOWED_ACTIONS)}。"
                "如果不需要动作，actions 返回空数组。"
            ),
        },
        {
            "role": "user",
            "content": json.dumps({"prompt": prompt, "snapshot": snapshot}, ensure_ascii=False),
        },
    ]
    parsed = complete_json(messages, temperature=0.35, max_tokens=2000, retries=1)
    reply = str(parsed.get("reply") or "").strip()
    actions = parsed.get("actions") if isinstance(parsed.get("actions"), list) else []
    normalized = []
    for item in actions:
        if not isinstance(item, dict):
            continue
        action_type = str(item.get("type") or "").strip()
        if action_type not in ALLOWED_ACTIONS:
            continue
        normalized.append(
            {
                "type": action_type,
                "params": item.get("params") if isinstance(item.get("params"), dict) else {},
                "reason": str(item.get("reason") or "").strip(),
                "ui_target": item.get("ui_target") if isinstance(item.get("ui_target"), dict) else {},
            }
        )
    return {"reply": reply, "actions": normalized, "snapshot": snapshot}


def execute_actions(
    db: Session,
    *,
    project_id: int,
    current_user: User,
    actions: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    results: list[dict[str, Any]] = []
    for action in actions:
        action_type = action["type"]
        params = action.get("params", {})
        result = {
            "type": action_type,
            "params": params,
            "reason": action.get("reason", ""),
            "ui_target": action.get("ui_target", {}),
            "result": {},
            "status": "ok",
        }
        try:
            if action_type == "navigate_page":
                result["result"] = {"message": "已生成导航指令", "path": params.get("path") or action.get("ui_target", {}).get("path", "")}
            elif action_type == "open_task":
                task = _require_task(db, project_id, _coerce_int(params.get("task_id")))
                result["result"] = {"message": "已定位任务", "task_id": task.id}
            elif action_type == "open_document":
                document = _require_document(db, project_id, _coerce_int(params.get("document_id")))
                result["result"] = {"message": "已定位文档", "document_id": document.id}
            elif action_type == "open_file":
                file_row = _require_file(db, project_id, _coerce_int(params.get("file_id")))
                result["result"] = {"message": "已定位文件", "file_id": file_row.id}
            elif action_type == "create_task":
                task = Task(
                    project_id=project_id,
                    title=str(params.get("title") or "AI 创建任务").strip(),
                    description=str(params.get("description") or "").strip(),
                    priority=str(params.get("priority") or "中").strip(),
                    assignee_id=_coerce_int(params.get("assignee_id")),
                    due_date=str(params.get("due_date") or "").strip(),
                    milestone_id=_coerce_int(params.get("milestone_id")),
                    status="TODO",
                    response_status="待接收" if _coerce_int(params.get("assignee_id")) else "未分配",
                    created_by=current_user.id,
                )
                db.add(task)
                db.flush()
                db.add(
                    Conversation(
                        project_id=project_id,
                        conversation_type="task",
                        related_task_id=task.id,
                        name=f"任务 #{task.id} {task.title}",
                        created_by=current_user.id,
                    )
                )
                db.add(
                    TaskActivity(
                        task_id=task.id,
                        project_id=project_id,
                        actor_id=current_user.id,
                        activity_type="created",
                        content=f"{current_user.display_name} 通过 AI 创建了任务 {task.title}",
                    )
                )
                db.commit()
                result["result"] = {"message": "任务已创建", "task_id": task.id, "title": task.title}
            elif action_type == "update_task":
                task = _require_task(db, project_id, _coerce_int(params.get("task_id")))
                for field in ["title", "description", "priority", "due_date"]:
                    if params.get(field) not in (None, ""):
                        setattr(task, field, str(params[field]).strip())
                if params.get("assignee_id") is not None:
                    task.assignee_id = _coerce_int(params.get("assignee_id"))
                if params.get("progress") is not None:
                    task.progress = max(0, min(100, _coerce_int(params.get("progress")) or 0))
                if params.get("status"):
                    task.status = str(params["status"]).strip()
                db.add(
                    TaskActivity(
                        task_id=task.id,
                        project_id=project_id,
                        actor_id=current_user.id,
                        activity_type="updated",
                        content=f"{current_user.display_name} 通过 AI 更新了任务 {task.title}",
                    )
                )
                db.commit()
                result["result"] = {"message": "任务已更新", "task_id": task.id}
            elif action_type == "accept_task":
                task = _require_task(db, project_id, _coerce_int(params.get("task_id")))
                task.status = "IN_PROGRESS"
                task.response_status = "处理中"
                db.add(
                    TaskActivity(
                        task_id=task.id,
                        project_id=project_id,
                        actor_id=current_user.id,
                        activity_type="accepted",
                        content=f"{current_user.display_name} 通过 AI 接收了任务 {task.title}",
                    )
                )
                db.commit()
                result["result"] = {"message": "任务已开始处理", "task_id": task.id}
            elif action_type == "block_task":
                task = _require_task(db, project_id, _coerce_int(params.get("task_id")))
                task.status = "BLOCKED"
                task.response_status = "阻塞中"
                task.blocker_reason = str(params.get("blocker_reason") or params.get("reason") or "待补充原因").strip()
                db.add(
                    TaskActivity(
                        task_id=task.id,
                        project_id=project_id,
                        actor_id=current_user.id,
                        activity_type="blocked",
                        content=f"{current_user.display_name} 通过 AI 标记任务阻塞：{task.blocker_reason}",
                    )
                )
                db.commit()
                result["result"] = {"message": "任务已标记阻塞", "task_id": task.id}
            elif action_type == "complete_task":
                task = _require_task(db, project_id, _coerce_int(params.get("task_id")))
                task.status = "DONE"
                task.response_status = "已完成"
                task.progress = 100
                db.add(
                    TaskActivity(
                        task_id=task.id,
                        project_id=project_id,
                        actor_id=current_user.id,
                        activity_type="completed",
                        content=f"{current_user.display_name} 通过 AI 完成了任务 {task.title}",
                    )
                )
                db.commit()
                result["result"] = {"message": "任务已完成", "task_id": task.id}
            elif action_type == "create_milestone":
                row = Milestone(
                    project_id=project_id,
                    name=str(params.get("name") or "AI 创建里程碑").strip(),
                    due_date=str(params.get("due_date") or "").strip(),
                    status=str(params.get("status") or "待开始").strip(),
                    order_index=db.query(Milestone).filter(Milestone.project_id == project_id).count(),
                )
                db.add(row)
                db.commit()
                db.refresh(row)
                result["result"] = {"message": "里程碑已创建", "milestone_id": row.id}
            elif action_type == "generate_weekly_report":
                report = generate_weekly_report(db, project_id)
                row = AIReportHistory(project_id=project_id, report_type="weekly", content=report)
                db.add(row)
                db.commit()
                result["result"] = {"message": "周报已生成", "content": report}
            elif action_type == "scan_risks":
                risk = RiskAlert(
                    project_id=project_id,
                    title=str(params.get("title") or "AI 风险扫描结果").strip(),
                    level=str(params.get("level") or "中").strip(),
                    score=float(params.get("score") or 60),
                    reason=str(params.get("reason") or "AI 结合当前项目数据识别到新的协作风险。").strip(),
                    suggestion=str(params.get("suggestion") or "请尽快确认责任人、截止时间和缓解动作。").strip(),
                    status="open",
                    risk_type=str(params.get("risk_type") or "协作风险").strip(),
                    due_at=str(params.get("due_at") or "").strip(),
                )
                db.add(risk)
                db.commit()
                db.refresh(risk)
                result["result"] = {"message": "风险扫描结果已写入", "risk_id": risk.id}
            elif action_type == "recalculate_contribution":
                payload = build_contribution_payload(db, project_id, prefer_llm=True, persist=True)
                db.commit()
                result["result"] = {"message": "贡献分析已重算", "analysis_mode": payload.get("summary", {}).get("analysis_mode")}
        except Exception as exc:
            db.rollback()
            result["status"] = "error"
            result["result"] = {"message": str(exc)}
        results.append(result)
    return results


def extract_file_text(storage_path: str, file_name: str) -> str:
    suffix = Path(file_name).suffix.lower()
    path = Path(storage_path)
    if not path.exists():
        raise ValueError("文件不存在或尚未落盘")
    if suffix in {".txt", ".md", ".json", ".py", ".ts", ".js", ".sql", ".csv", ".html", ".css"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        return "\n".join((page.extract_text() or "") for page in reader.pages)
    if suffix == ".docx":
        from docx import Document as DocxDocument

        doc = DocxDocument(str(path))
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(str(path))
        chunks = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
        return "\n".join(chunks)
    raise ValueError(f"暂不支持分析 {suffix or '该类型'} 文件")


def _messages_page_payload(db: Session, project_id: int, selected_conversation_id: int | None) -> dict[str, Any]:
    conversations = db.query(Conversation).filter(Conversation.project_id == project_id).order_by(Conversation.created_at.desc()).all()
    selected = next((item for item in conversations if item.id == selected_conversation_id), conversations[0] if conversations else None)
    users = build_user_map(db)
    payload = {
        "conversations": [
            {
                "id": item.id,
                "name": item.name,
                "conversation_type": item.conversation_type,
            }
            for item in conversations[:12]
        ]
    }
    if selected:
        messages = (
            db.query(Message)
            .filter(Message.conversation_id == selected.id)
            .order_by(Message.created_at.desc())
            .limit(8)
            .all()
        )
        payload["selected_conversation"] = {
            "id": selected.id,
            "name": selected.name,
            "messages": [serialize_message(item, users) for item in reversed(messages)],
        }
    return payload


def _task_brief(task: Task | None, users: dict[int, User]) -> dict[str, Any] | None:
    if not task:
        return None
    return serialize_task(task, users)


def _document_brief(document: Document | None, *, include_content: bool = False) -> dict[str, Any] | None:
    if not document:
        return None
    return {
        "id": document.id,
        "title": document.title,
        "updated_at": str(document.updated_at),
        "related_task_id": document.related_task_id,
        "content": document.content[:4000] if include_content else "",
    }


def _file_brief(file_row: FileResource | None, *, include_description: bool = False) -> dict[str, Any] | None:
    if not file_row:
        return None
    return {
        "id": file_row.id,
        "name": file_row.name,
        "category": file_row.category,
        "review_status": file_row.review_status,
        "size_label": file_row.size_label,
        "description": file_row.description if include_description else "",
    }


def _require_task(db: Session, project_id: int, task_id: int | None) -> Task:
    if not task_id:
        raise ValueError("缺少 task_id")
    row = db.query(Task).filter(Task.project_id == project_id, Task.id == task_id, Task.archived == False).first()  # noqa: E712
    if not row:
        raise ValueError("任务不存在")
    return row


def _require_document(db: Session, project_id: int, document_id: int | None) -> Document:
    if not document_id:
        raise ValueError("缺少 document_id")
    row = db.query(Document).filter(Document.project_id == project_id, Document.id == document_id, Document.is_deleted == False).first()  # noqa: E712
    if not row:
        raise ValueError("文档不存在")
    return row


def _require_file(db: Session, project_id: int, file_id: int | None) -> FileResource:
    if not file_id:
        raise ValueError("缺少 file_id")
    row = db.query(FileResource).filter(FileResource.project_id == project_id, FileResource.id == file_id).first()
    if not row:
        raise ValueError("文件不存在")
    return row


def _coerce_int(value: Any) -> int | None:
    try:
        if value in (None, ""):
            return None
        return int(value)
    except (TypeError, ValueError):
        return None
